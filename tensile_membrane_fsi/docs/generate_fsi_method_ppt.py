#!/usr/bin/env python3
"""Generate a step-by-step PowerPoint on the tensile-membrane FSI method."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(__file__).resolve().parent / "Tensile_Membrane_FSI_Method.pptx"

# Visual system — avoid purple / cream AI defaults
NAVY = RGBColor(0x0B, 0x2E, 0x3D)
TEAL = RGBColor(0x1F, 0x7A, 0x6B)
SLATE = RGBColor(0x2C, 0x3E, 0x50)
MUTED = RGBColor(0x5A, 0x6A, 0x75)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF6, 0xF7)
ACCENT = RGBColor(0xC4, 0x5C, 0x26)  # burnt orange, sparingly


def _set_run(run, size=18, bold=False, color=SLATE, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _fill_solid(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_bg(slide, prs, color=LIGHT):
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0),
        Inches(0),
        prs.slide_width,
        prs.slide_height,
    )
    _fill_solid(shape, color)
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_header_bar(slide, prs, title: str, step: str | None = None):
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.95)
    )
    _fill_solid(bar, NAVY)
    accent = slide.shapes.add_shape(
        1, Inches(0), Inches(0.95), prs.slide_width, Inches(0.06)
    )
    _fill_solid(accent, TEAL)

    box = slide.shapes.add_textbox(
        Inches(0.45), Inches(0.22), Inches(11.5), Inches(0.55)
    )
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=26, bold=True, color=WHITE, font="Georgia")

    if step:
        tag = slide.shapes.add_textbox(
            Inches(10.4), Inches(0.28), Inches(2.3), Inches(0.4)
        )
        tp = tag.text_frame.paragraphs[0]
        tp.alignment = PP_ALIGN.RIGHT
        r = tp.add_run()
        r.text = step
        _set_run(r, size=12, bold=True, color=RGBColor(0xA8, 0xD5, 0xCC))


def add_footer(slide, prs, page: int, total: int):
    foot = slide.shapes.add_textbox(
        Inches(0.45), Inches(7.05), Inches(11.5), Inches(0.3)
    )
    p = foot.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"Tensile Membrane FSI  ·  Partitioned serial-staggered method  ·  {page}/{total}"
    _set_run(r, size=10, color=MUTED)


def add_bullets(slide, left, top, width, height, items, size=16, color=SLATE):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = item[1] if isinstance(item, tuple) else 0
        text = item[0] if isinstance(item, tuple) else item
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = ("• " if p.level == 0 else "– ") + text
        _set_run(run, size=size - (2 if p.level else 0), color=color)
    return box


def add_card(slide, left, top, width, height, title, body_lines, title_color=TEAL):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xD0, 0xD8, 0xDC)

    tbox = slide.shapes.add_textbox(
        left + Inches(0.18), top + Inches(0.12), width - Inches(0.3), Inches(0.35)
    )
    p = tbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run(r, size=15, bold=True, color=title_color, font="Georgia")

    bbox = slide.shapes.add_textbox(
        left + Inches(0.18),
        top + Inches(0.48),
        width - Inches(0.3),
        height - Inches(0.6),
    )
    tf = bbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        _set_run(run, size=12, color=SLATE)


def add_equation_box(slide, left, top, width, height, lines):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    tbox = slide.shapes.add_textbox(
        left + Inches(0.25), top + Inches(0.15), width - Inches(0.4), height - Inches(0.2)
    )
    tf = tbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = line
        _set_run(r, size=16, bold=True, color=WHITE, font="Cambria Math")


def title_slide(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs, NAVY)
    band = s.shapes.add_shape(1, Inches(0), Inches(5.6), prs.slide_width, Inches(1.9))
    _fill_solid(band, TEAL)

    t = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(11.5), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Tensile Membrane FSI"
    _set_run(r, size=44, bold=True, color=WHITE, font="Georgia")

    st = s.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(11.5), Inches(1.0))
    p = st.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Partitioned Serial-Staggered Method — Step by Step"
    _set_run(r, size=22, color=RGBColor(0xB8, 0xD9, 0xD2), font="Calibri")

    f = s.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(11.5), Inches(0.9))
    p = f.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = (
        "Fluid–structure interaction of prestressed fabric membranes\n"
        "in incompressible flow  ·  tensile_membrane_fsi"
    )
    _set_run(r, size=14, color=WHITE)


def agenda_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "Agenda — What this deck covers")
    items = [
        "Problem & modelling goal",
        "Partitioned vs monolithic FSI",
        "Code architecture & modules",
        "Setup: mesh, material, prestress, fluid grid",
        "One coupled time step — every sub-step",
        "Membrane dynamics (CST FEM + explicit integrator)",
        "Fluid solver (PISO-like + LES + immersed boundary)",
        "Load transfer modes",
        "Stabilisation knobs (damping, mass_scale, under-relaxation)",
        "Outputs, flutter case, and static DOLFINx contrast",
    ]
    add_bullets(s, Inches(0.7), Inches(1.35), Inches(11), Inches(5.4), items, size=17)
    add_footer(s, prs, page, total)


def problem_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "1. Problem statement", "Part A")
    add_card(
        s,
        Inches(0.5),
        Inches(1.3),
        Inches(5.8),
        Inches(5.2),
        "Physics",
        [
            "A thin prestressed tensile membrane",
            "(sail / canopy / fabric roof) sits in a",
            "3D air channel.",
            "",
            "Flow exerts unsteady pressure → membrane",
            "deforms and moves.",
            "",
            "Membrane motion changes the flow",
            "(immersed boundary) → two-way FSI.",
            "",
            "Goal: capture transient interaction,",
            "including flutter when damping is low.",
        ],
    )
    add_card(
        s,
        Inches(6.6),
        Inches(1.3),
        Inches(5.9),
        Inches(5.2),
        "Numerical approach in this repo",
        [
            "Membrane: CST triangles, prestress +",
            "plane-stress elasticity, explicit CD.",
            "",
            "Fluid: incompressible NS on Cartesian",
            "grid, fractional-step / PISO-like,",
            "optional Smagorinsky LES.",
            "",
            "Coupling: partitioned serial staggered",
            "(not monolithic).",
            "",
            "IB: thin-band mask of fluid cells on the",
            "membrane surface.",
        ],
    )
    add_footer(s, prs, page, total)


def partitioned_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "2. Partitioned vs monolithic", "Part A")
    add_card(
        s,
        Inches(0.5),
        Inches(1.3),
        Inches(5.8),
        Inches(5.2),
        "This code = PARTITIONED",
        [
            "Fluid and structure are separate solvers.",
            "",
            "Each time step:",
            "  1) advance fluid with current IB",
            "  2) transfer loads to membrane",
            "  3) advance membrane",
            "  4) update IB from new position",
            "",
            "Optional under-relaxed sub-iterations",
            "tighten the interface residual.",
            "",
            "Config: scheme: serial_staggered",
        ],
        title_color=TEAL,
    )
    add_card(
        s,
        Inches(6.6),
        Inches(1.3),
        Inches(5.9),
        Inches(5.2),
        "Monolithic (NOT used here)",
        [
            "Fluid + structure + interface assembled",
            "into one joint algebraic system.",
            "",
            "Solved together each step / Newton",
            "iteration.",
            "",
            "Stronger coupling / often more stable,",
            "but larger, more complex Jacobian and",
            "harder to reuse separate solvers.",
            "",
            "Your DOLFINx SNES script is membrane-",
            "only static analysis — neither FSI style.",
        ],
        title_color=ACCENT,
    )
    add_footer(s, prs, page, total)


def architecture_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "3. Repository architecture", "Part A")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.25),
        Inches(12),
        Inches(5.5),
        [
            "main.py — load YAML, build FSISimulation, run loop, save outputs",
            "config/default.yaml | flutter.yaml — physics & numerics",
            "src/membrane/ — geometry (CST mesh), materials, prestress / form-finding, explicit solver",
            "src/fluid/ — Cartesian grid, PISO-like NS, Smagorinsky LES",
            "src/fsi/ — coupling driver, load transfer, immersed-boundary update",
            "src/utils/ — YAML I/O, NPZ/VTK writers, plots, GIF helpers",
            "examples/ — membrane-only, coarse FSI, flutter GIF, Gmsh generator",
            "tests/ — unit / smoke checks for membrane and coupled step",
        ],
        size=16,
    )
    add_footer(s, prs, page, total)


def setup_overview_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "4. Setup pipeline (before time stepping)", "Part B")
    steps = [
        ("A", "Build rectangular CST membrane mesh + fixed-edge BCs"),
        ("B", "Apply slight analytical sag so flow sees a curved sail"),
        ("C", "Create MembraneMaterial (E, ν, thickness, prestress)"),
        ("D", "Dynamic-relaxation prestress (heavily damped) → equilibrium"),
        ("E", "Restore user transient damping from config"),
        ("F", "Build FluidGrid + FluidSolver (ρ, ν, U_inlet, LES flags)"),
        ("G", "Initial immersed-boundary mask from membrane state"),
        ("H", "Read FSI knobs: α, max_subiters, load_mode, load_scale"),
    ]
    y = 1.25
    for tag, text in steps:
        chip = s.shapes.add_shape(1, Inches(0.55), Inches(y), Inches(0.45), Inches(0.42))
        _fill_solid(chip, TEAL)
        tb = s.shapes.add_textbox(Inches(0.55), Inches(y + 0.05), Inches(0.45), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = tag
        _set_run(r, size=14, bold=True, color=WHITE)
        txt = s.shapes.add_textbox(Inches(1.2), Inches(y + 0.02), Inches(11), Inches(0.4))
        p = txt.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = text
        _set_run(r, size=16, color=SLATE)
        y += 0.65
    add_footer(s, prs, page, total)


def membrane_mesh_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "5. Membrane mesh & BCs", "Part B · Setup A–B")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5.3),
        [
            "Rectangular patch in a reference plane; discretised into constant-strain triangles (CST)",
            "Config: length, width, nx, ny → node/element connectivity",
            "Placement in fluid: membrane_x0, membrane_y0, membrane_z0",
            "fixed_edges (e.g. left/right/bottom/top) → Dirichlet nodes held at reference",
            "initial_sag_shape: sinusoidal z-displacement so the sail is not perfectly flat",
            "Why sag? A flat zero-curvature start has weak geometric stiffness and is a poor IB seed",
            "Geometry helpers recompute element areas and normals as the membrane moves",
        ],
        size=16,
    )
    add_footer(s, prs, page, total)


def material_prestress_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "6. Material & prestress form-finding", "Part B · Setup C–E")
    add_equation_box(
        s,
        Inches(0.55),
        Inches(1.25),
        Inches(12.0),
        Inches(1.15),
        [
            "Plane-stress elasticity + isotropic prestress resultant N_pre",
            "Form-finding: heavily damped dynamics → near-equilibrium before FSI",
        ],
    )
    add_bullets(
        s,
        Inches(0.6),
        Inches(2.7),
        Inches(12),
        Inches(4.0),
        [
            "MembraneMaterial: Young’s modulus E, Poisson ν, thickness h, prestress",
            "apply_isotropic_prestress runs MembraneSolver with high damping + mass_scale",
            "Many small explicit steps; velocities periodically halved, then zeroed",
            "Purpose: settle fabric under prestress + gravity (dynamic relaxation)",
            "After form-finding, transient damping is restored from config (may be ~0 for flutter)",
            "This prestress damping is NOT the same as wanting viscous damping in the FSI run",
        ],
        size=15,
    )
    add_footer(s, prs, page, total)


def fluid_setup_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "7. Fluid grid & solver setup", "Part B · Setup F–G")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5.3),
        [
            "Uniform Cartesian box: domain L × W × H, resolution nx × ny × nz",
            "State fields: velocity (u,v,w), pressure p, effective viscosity ν_eff",
            "Inlet profile tapered near walls to avoid corner spikes",
            "Optional gust: sinusoidal streamwise + vertical inlet fluctuation",
            "LES: Smagorinsky SGS adds eddy viscosity when enabled",
            "Initial IB: membrane_cell_mask marks a thin band of cells as solid",
            "Solid velocity in masked cells interpolated from nearest membrane nodes",
        ],
        size=16,
    )
    add_footer(s, prs, page, total)


def timestep_overview_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "8. One FSI time step — overview", "Part C")
    add_equation_box(
        s,
        Inches(0.55),
        Inches(1.25),
        Inches(12.0),
        Inches(0.85),
        ["Serial staggered loop inside FSISimulation.step()  (repeat ≤ max_subiters)"],
    )
    items = [
        "1. Update immersed boundary from current membrane x, v",
        "2. Advance fluid one Δt (advection–diffusion → pressure Poisson → projection)",
        "3. Compute surface loads → nodal forces (× load_scale)",
        "4. Under-relax forces if sub-iteration > 0: f ← α f_new + (1−α) f_old",
        "5. Set membrane external forces; advance membrane with n_sub substeps",
        "6. Measure interface residual from displacement change; break if < tol",
        "7. After sub-iters: t += Δt; log disp, KE, CFL, residual, |p|_max",
    ]
    add_bullets(s, Inches(0.7), Inches(2.35), Inches(11.5), Inches(4.3), items, size=16)
    add_footer(s, prs, page, total)


def step_ib_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "9. Step 1 — Update immersed boundary", "Part C · Detail")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5.3),
        [
            "Function: update_immersed_boundary(...) in src/fsi/mesh_update.py",
            "Build boolean mask of fluid cells near the membrane (thickness_cells ≈ 1.5)",
            "For each masked cell, find nearest membrane node (in x–y) and copy its velocity",
            "Pass mask + solid velocity (u_s, v_s, w_s) into FluidSolver.set_immersed_boundary",
            "During fluid BC application, masked cells are forced to solid velocity",
            "This is how structure → fluid coupling is enforced (kinematic condition)",
            "As the membrane moves each step / sub-iter, the mask is rebuilt",
        ],
        size=16,
    )
    add_footer(s, prs, page, total)


def step_fluid_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "10. Step 2 — Advance the fluid", "Part C · Detail")
    add_equation_box(
        s,
        Inches(0.55),
        Inches(1.2),
        Inches(12.0),
        Inches(1.35),
        [
            "Fractional-step / PISO-like incompressible NS",
            "u* ← advect–diffuse(u)   →   ∇²φ = (ρ/Δt) ∇·u*   →   u ← u* − (Δt/ρ) ∇φ",
        ],
    )
    add_bullets(
        s,
        Inches(0.6),
        Inches(2.8),
        Inches(12),
        Inches(3.8),
        [
            "Convection: first-order upwind; diffusion: Laplacian with ν_eff",
            "ν_eff = ν_molecular (+ Smagorinsky eddy viscosity if LES on)",
            "Pressure Poisson with Neumann walls; one cell pinned for uniqueness",
            "BCs: inlet (profile ± gust), no-slip side walls, zero-gradient outlet, IB force",
            "Sanitizer clips NaN / runaway velocities (u_clip) on coarse teaching grids",
            "CFL monitored after the coupled step for diagnostics",
        ],
        size=15,
    )
    add_footer(s, prs, page, total)


def step_loads_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "11. Step 3 — Load transfer (fluid → structure)", "Part C · Detail")
    add_card(
        s,
        Inches(0.4),
        Inches(1.25),
        Inches(4.0),
        Inches(5.2),
        "dynamic_pressure",
        [
            "Default teaching model.",
            "",
            "Sample fluid vel & p at",
            "element centroids.",
            "",
            "Δp ≈ ½ ρ U_n |U_n|",
            "+ 0.3 (p − p∞)",
            "",
            "Clip to ±5 q_ref.",
            "",
            "Incidence / Bernoulli-like;",
            "good for angled loading.",
        ],
    )
    add_card(
        s,
        Inches(4.55),
        Inches(1.25),
        Inches(4.0),
        Inches(5.2),
        "pressure_jump",
        [
            "Used in flutter.yaml.",
            "",
            "Sample p at",
            "centroid ± offset·n",
            "(outside IB band).",
            "",
            "Δp = p(−n) − p(+n)",
            "along +normal.",
            "",
            "Captures unsteady lift",
            "even when membrane is",
            "aligned with the flow.",
        ],
    )
    add_card(
        s,
        Inches(8.7),
        Inches(1.25),
        Inches(4.0),
        Inches(5.2),
        "interpolated_field",
        [
            "Use interpolated fluid",
            "pressure directly as",
            "surface pressure.",
            "",
            "Simplest mapping;",
            "depends on how well",
            "IB pressure is defined",
            "near the membrane.",
            "",
            "All modes then scatter",
            "Fe = (p A / 3) n to the",
            "three triangle nodes.",
        ],
    )
    add_footer(s, prs, page, total)


def step_relax_membrane_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "12. Steps 4–6 — Relax, integrate membrane, residual", "Part C · Detail")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.25),
        Inches(12),
        Inches(5.5),
        [
            "load_scale multiplies nodal forces (stabilises coarse / aggressive FSI)",
            "Under-relaxation (sub-iter > 0): f = α f_new + (1−α) f_old  — reduces load oscillation",
            "membrane.set_external_forces(f_use)",
            "n_mem = required_substeps(Δt) from membrane critical Δt (wave-speed CFL), capped at 200",
            "membrane.step(Δt, n_sub=n_mem) advances structure with smaller structural substeps",
            "Residual ≈ ‖Δx‖ / ‖x_prev‖ between sub-iterations; stop if < residual_tol",
            "max_subiters = 1 → pure serial stagger (weak coupling); >1 → iterated weak coupling",
        ],
        size=15,
    )
    add_footer(s, prs, page, total)


def membrane_dynamics_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "13. Membrane dynamics — residual & integrator", "Part D")
    add_equation_box(
        s,
        Inches(0.55),
        Inches(1.2),
        Inches(12.0),
        Inches(1.5),
        [
            "M a = f_ext + f_gravity − f_int(x) − c v",
            "Explicit central difference on free DOFs; fixed edges held at x₀",
        ],
    )
    add_bullets(
        s,
        Inches(0.6),
        Inches(3.0),
        Inches(12),
        Inches(3.6),
        [
            "f_int = material (Green–Lagrange / plane-stress via CST B matrix) + geometric prestress",
            "Lumped mass × mass_scale for artificial inertia (larger stable fluid Δt)",
            "Gravity uses unscaled physical mass so weight is not artificially amplified",
            "c = damping (Rayleigh-like viscous); set ~0 for persistent flutter",
            "Safety: clip runaway nodal displacements; zero velocity on clipped nodes",
            "critical_dt from max(material, prestress) wave speed and min edge length",
        ],
        size=15,
    )
    add_footer(s, prs, page, total)


def membrane_internal_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "14. Membrane internal forces (per triangle)", "Part D")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.25),
        Inches(12),
        Inches(5.5),
        [
            "Build local orthonormal frame (e1, e2, n) from reference triangle X₀",
            "Project displacements into the local plane → 2D CST B matrix",
            "Strain ε = B u_loc  (Voigt); stress σ = C ε; membrane resultant N = h σ",
            "Material nodal forces from Bᵀ (N A); map back to 3D via e1, e2",
            "Add geometric / prestress edge forces from isotropic N_pre",
            "Assemble to global f_int; residual subtracts f_int from loads + gravity",
            "Contrast: DOLFINx static script derives residual from total potential energy + SNES",
        ],
        size=15,
    )
    add_footer(s, prs, page, total)


def damping_mass_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "15. Damping & mass_scale — why they exist", "Part D")
    add_card(
        s,
        Inches(0.5),
        Inches(1.3),
        Inches(5.8),
        Inches(5.2),
        "damping c",
        [
            "Appears only in dynamics: −c v",
            "",
            "Dissipates KE; kills mesh-scale",
            "ringing and noisy FSI loads.",
            "",
            "NOT required by static load analysis",
            "(energy / SNES equilibrium).",
            "",
            "default.yaml: c ≈ 80 (quiet demo)",
            "flutter.yaml: c ≈ 0.02 (oscillations live)",
            "",
            "You may set c = 0 if you accept",
            "undamped / more fragile transients.",
        ],
    )
    add_card(
        s,
        Inches(6.6),
        Inches(1.3),
        Inches(5.9),
        Inches(5.2),
        "mass_scale",
        [
            "Multiplies lumped nodal mass for",
            "the explicit integrator.",
            "",
            "Slows structural response → larger",
            "stable Δt when coupled to fluid.",
            "",
            "Does NOT dissipate energy.",
            "",
            "Gravity uses unscaled mass so the",
            "sail weight stays physical.",
            "",
            "Trade-off: artificial inertia vs",
            "true membrane dynamics.",
        ],
    )
    add_footer(s, prs, page, total)


def les_ib_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "16. LES & immersed-boundary details", "Part E")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5.3),
        [
            "Smagorinsky: ν_t ~ (C_s Δ)² |S|; ν_eff = ν + ν_t when les.enabled",
            "Teaching grids often raise molecular ν and may disable LES (--quick)",
            "IB is a thin volumetric band, not a sharp zero-thickness interface",
            "pressure_jump offset ≈ 3 Δz samples outside that band for Δp",
            "Kinematic coupling: fluid velocity in band = membrane velocity",
            "Dynamic coupling: fluid pressure / dynamic pressure → nodal forces",
            "No ALE mesh morphing of the fluid grid — Cartesian + mask only",
        ],
        size=16,
    )
    add_footer(s, prs, page, total)


def run_loop_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "17. Outer run loop & outputs", "Part E")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.25),
        Inches(12),
        Inches(5.5),
        [
            "n_steps = ceil(t_end / dt); each iteration calls FSISimulation.step()",
            "Callback every save_interval: print diagnostics + write artifacts",
            "snapshot_XXXXXX.npz — membrane nodes + fluid fields",
            "membrane_XXXXXX.vtk — open in ParaView",
            "slice_XXXXXX.png — mid-plane |u| with membrane outline",
            "history.csv / history.png — displacement, KE, CFL, FSI residual",
            "membrane_flutter.gif — from examples/run_flutter_gif.py",
            "Entry points: python main.py [-c config] [--quick]; examples/*",
        ],
        size=15,
    )
    add_footer(s, prs, page, total)


def config_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "18. Key configuration knobs", "Part E")
    add_card(
        s,
        Inches(0.4),
        Inches(1.25),
        Inches(4.0),
        Inches(5.2),
        "membrane.*",
        [
            "length, width, nx, ny",
            "thickness, density",
            "youngs_modulus, poisson",
            "prestress",
            "damping, mass_scale",
            "fixed_edges",
        ],
    )
    add_card(
        s,
        Inches(4.55),
        Inches(1.25),
        Inches(4.0),
        Inches(5.2),
        "fluid.* / les.*",
        [
            "domain L,W,H + nx,ny,nz",
            "rho, nu, U_inlet",
            "membrane_x0/y0/z0",
            "gust_amp, gust_freq",
            "u_clip (sanitizer)",
            "les.enabled, Cs",
        ],
    )
    add_card(
        s,
        Inches(8.7),
        Inches(1.25),
        Inches(4.0),
        Inches(5.2),
        "fsi.* / time.*",
        [
            "scheme: serial_staggered",
            "under_relaxation α",
            "max_subiters",
            "residual_tol",
            "load_scale, load_mode",
            "dt, t_end, cfl_max",
        ],
    )
    add_footer(s, prs, page, total)


def flutter_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "19. Flutter demo path", "Part E")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5.3),
        [
            "config/flutter.yaml: light fabric, soft prestress, near-zero structural damping",
            "Gusty inflow to seed unsteady lift",
            "load_mode: pressure_jump for Δp across the sheet",
            "Clamped on all four edges; oscillations intended to persist",
            "examples/run_flutter_gif.py writes output/flutter/membrane_flutter.gif",
            "Committed artifacts: NPZ/VTK snapshots, slices, history, GIF",
            "Shows the partitioned method can sustain dynamic FSI without heavy damping",
        ],
        size=16,
    )
    add_footer(s, prs, page, total)


def dolfinx_contrast_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "20. Contrast: static DOLFINx load analysis", "Part F")
    add_card(
        s,
        Inches(0.5),
        Inches(1.3),
        Inches(5.8),
        Inches(5.2),
        "DOLFINx + PETSc SNES (your script)",
        [
            "Static nonlinear membrane only",
            "Total potential energy Π(u)",
            "Residual = δΠ; Jacobian = dF/du",
            "Point load with load factor stepping",
            "Small u₃ nudge for geometric stiffness",
            "No mass, no damping, no fluid",
            "Equilibrium: F_int(u) = F_ext",
            "MUMPS / Newton line search",
        ],
    )
    add_card(
        s,
        Inches(6.6),
        Inches(1.3),
        Inches(5.9),
        Inches(5.2),
        "This repo FSI path",
        [
            "Transient partitioned FSI",
            "Separate fluid + structure solvers",
            "Explicit CD structure + PISO fluid",
            "Interface loads each Δt",
            "Optional viscous damping",
            "Artificial mass scaling common",
            "IB kinematic coupling",
            "Teaching / prototyping fidelity",
        ],
    )
    add_footer(s, prs, page, total)


def limits_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "21. Limits & practical tips", "Part F")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5.3),
        [
            "Cartesian IB fluid ≠ production unstructured LES / OpenFOAM pipeline",
            "Explicit membrane Δt limited by wave speed — lower dt or E / prestress if unstable",
            "Dynamic-pressure load model is approximate; swap load_transfer.py for fidelity",
            "Raised ν and load_scale are stability aids on coarse grids",
            "If you remove damping: expect more ringing; reduce dt / mass_scale / load_scale first",
            "For static design loads, prefer energy-based SNES (DOLFINx) over this transient FSI",
            "Gmsh helper exists for external unstructured coupling experiments",
        ],
        size=15,
    )
    add_footer(s, prs, page, total)


def flowchart_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "22. End-to-end flowchart", "Summary")
    boxes = [
        (0.5, "YAML\nconfig"),
        (2.5, "Build mesh\n+ material"),
        (4.5, "Prestress\nrelax"),
        (6.5, "Fluid +\nIB init"),
        (8.5, "Time loop\nFSI.step"),
        (10.5, "Write\noutputs"),
    ]
    for x, label in boxes:
        shp = s.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(1.8), Inches(1.2))
        _fill_solid(shp, TEAL if x != 8.5 else ACCENT)
        tb = s.shapes.add_textbox(Inches(x), Inches(2.25), Inches(1.8), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        _set_run(r, size=13, bold=True, color=WHITE)
        if x < 10.5:
            arr = s.shapes.add_textbox(Inches(x + 1.75), Inches(2.35), Inches(0.4), Inches(0.4))
            p = arr.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = "→"
            _set_run(r, size=20, bold=True, color=NAVY)

    inner = [
        "Inside FSI.step (sub-iters):  IB update → fluid.step → loads → under-relax → membrane.step → residual check",
        "Membrane subcycles with n_sub from structural CFL while fluid advances with global Δt",
        "History accumulates time, max displacement, kinetic energy, CFL, coupling residual",
    ]
    add_bullets(s, Inches(0.6), Inches(3.6), Inches(12), Inches(2.8), inner, size=15)
    add_footer(s, prs, page, total)


def takeaways_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header_bar(s, prs, "23. Key takeaways", "Summary")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(12),
        Inches(5.3),
        [
            "Coupling is partitioned & serial staggered — not monolithic",
            "Two-way link: IB (structure→fluid) + pressure loads (fluid→structure)",
            "Membrane: CST + prestress + explicit CD; fluid: PISO-like + optional LES",
            "Damping is optional physics/numerics for dynamics; unused in static SNES analysis",
            "mass_scale, load_scale, α, sub-iters are the main coupling stabilisers",
            "Choose load_mode for the physics you care about (incidence vs Δp flutter)",
            "Use default.yaml for quiet demos; flutter.yaml for near-undamped oscillation",
        ],
        size=16,
    )
    add_footer(s, prs, page, total)


def closing_slide(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs, NAVY)
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.0))
    p = t.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Questions & next steps"
    _set_run(r, size=40, bold=True, color=WHITE, font="Georgia")

    st = s.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.5), Inches(2.0))
    tf = st.text_frame
    lines = [
        "Explore: config/default.yaml · config/flutter.yaml · src/fsi/coupling.py",
        "Run: python main.py --quick   |   python examples/run_flutter_gif.py",
        "Static loads: use DOLFINx energy / SNES path (no damping required)",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = line
        _set_run(r, size=16, color=RGBColor(0xB8, 0xD9, 0xD2))


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        title_slide,
        agenda_slide,
        problem_slide,
        partitioned_slide,
        architecture_slide,
        setup_overview_slide,
        membrane_mesh_slide,
        material_prestress_slide,
        fluid_setup_slide,
        timestep_overview_slide,
        step_ib_slide,
        step_fluid_slide,
        step_loads_slide,
        step_relax_membrane_slide,
        membrane_dynamics_slide,
        membrane_internal_slide,
        damping_mass_slide,
        les_ib_slide,
        run_loop_slide,
        config_slide,
        flutter_slide,
        dolfinx_contrast_slide,
        limits_slide,
        flowchart_slide,
        takeaways_slide,
        closing_slide,
    ]
    total = len(builders)
    # title has no page footer numbering in same style; still count in total
    for i, fn in enumerate(builders, start=1):
        if fn is title_slide:
            fn(prs, total)
        else:
            fn(prs, total, i)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT} ({total} slides)")


if __name__ == "__main__":
    build()
